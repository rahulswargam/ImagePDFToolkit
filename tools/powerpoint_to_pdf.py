import os

from tools.image_resizer import get_output_folder
from tools.office_com import is_office_app_installed

_PP_SAVE_AS_PDF = 32


def _unique_pdf_path(output_folder, base_name, suffix):
    output_path = os.path.join(output_folder, f"{base_name}{suffix}.pdf")
    counter = 1
    while os.path.exists(output_path):
        output_path = os.path.join(output_folder, f"{base_name}{suffix}_{counter}.pdf")
        counter += 1
    return output_path


def _convert_via_com(input_path, output_path):
    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()
    powerpoint = None
    try:
        # PowerPoint's COM automation doesn't reliably support `Visible = False`
        # the way Word/Excel do — WithWindow=False on Open is the correct way
        # to suppress the visible flash for this app specifically.
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(os.path.abspath(input_path), WithWindow=False)
        try:
            presentation.SaveAs(output_path, _PP_SAVE_AS_PDF)
        finally:
            presentation.Close()
    finally:
        if powerpoint is not None:
            powerpoint.Quit()
        pythoncom.CoUninitialize()


def _convert_via_fallback(input_path, output_path):
    if input_path.lower().endswith(".ppt"):
        raise ValueError("Legacy .ppt files need Microsoft PowerPoint installed to convert.")

    import pptx
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

    presentation = pptx.Presentation(input_path)
    styles = getSampleStyleSheet()
    story = []

    for slide in presentation.slides:
        has_text = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)
                if text.strip():
                    story.append(Paragraph(text, styles["Normal"]))
                    story.append(Spacer(1, 4))
                    has_text = True
        if not has_text:
            story.append(Paragraph("(slide has no extractable text)", styles["Italic"]))
        story.append(PageBreak())

    if not story:
        raise ValueError("This presentation has no slides.")

    SimpleDocTemplate(output_path, pagesize=LETTER).build(story)


def convert_powerpoint_to_pdf(input_path):
    """
    Convert a PowerPoint presentation into a PDF. Uses Microsoft PowerPoint
    (COM automation) when it's installed, for exact fidelity. Falls back to
    plain per-slide text extraction (no images, layout, or design — text
    only, one page break per slide) when PowerPoint isn't available or the
    COM call fails, since python-pptx has no rendering engine to fall back
    on more faithfully than that.

    Returns:
        (output_path, method) where method is "com" or "fallback"
    """

    input_path = str(input_path)
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    output_path = _unique_pdf_path(get_output_folder(), base_name, "")

    if is_office_app_installed("powerpoint"):
        try:
            _convert_via_com(input_path, output_path)
            return output_path, "com"
        except Exception:
            pass

    _convert_via_fallback(input_path, output_path)
    return output_path, "fallback"
